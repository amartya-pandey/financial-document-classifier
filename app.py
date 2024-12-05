import os
import dill
import pandas as pd
# import time
from flask import Flask, request, render_template
import chardet  # To detect file encoding
from io import BytesIO
from PyPDF2 import PdfReader  # For PDF files
from pptx import Presentation
# from retrain_model import main as ret_mod


app = Flask(__name__)

# Defines the categories
def get_subdirectories(path):
    return [d for d in os.listdir(path) if os.path.isdir(os.path.join(path, d))]
categories = get_subdirectories(r'data')


def extract_text_from_pptx(file):
    content = ""
    presentation = Presentation(file) 
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content += shape.text + "\n"
    return content.strip()  



categories = [
            'Consolidated statement of cash flows',
            'Statement of operations',
            'Statements of Financial Position', 
            'Statement of changes in equity',
            'Note to financial statements'  
            ]

import warnings
from sklearn.exceptions import InconsistentVersionWarning
warnings.filterwarnings("ignore", category=InconsistentVersionWarning)

# Load the model and vectorizer
with open('model.dill', 'rb') as model_file:
    model = dill.load(model_file)

with open('vectorizer.dill', 'rb') as vectorizer_file:
    vectorizer = dill.load(vectorizer_file)
def read_file_content(file):
    raw_data = file.read()
    result = chardet.detect(raw_data)
    encoding = result['encoding'] if result['encoding'] is not None else 'utf-8'

    file.seek(0)

    try:
        content = raw_data.decode(encoding)
    except UnicodeDecodeError:
        if file.filename.endswith('.pdf'):
            reader = PdfReader(BytesIO(raw_data))
            content = ""
            for page in reader.pages:
                content += page.extract_text() or ""
        elif file.filename.endswith('.pptx'):
            content = extract_text_from_pptx(file)
        else:
            raise ValueError("Unsupported file format or encoding.")
    
    return content


def store_feedback(content, predicted_category):
    """Store the feedback in feedback.csv."""
    feedback_path = 'feedback.csv'
    
    # Prepare the feedback data
    feedback_data = {'text': [content], 'label': [predicted_category]}
    feedback_df = pd.DataFrame(feedback_data)
    
    # Append feedback to the CSV file
    if os.path.exists(feedback_path):
        feedback_df.to_csv(feedback_path, mode='a', header=False, index=False)
    else:
        feedback_df.to_csv(feedback_path, index=False)


@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        file = request.files['file']
        if file:
            filename = file.filename  # Get the filename

            try:
                content = read_file_content(file)
                vectorized_content = vectorizer.transform([content])
                prediction_index = model.predict(vectorized_content)[0]
                predicted_category = categories[prediction_index]
                store_feedback(content, predicted_category)
                return render_template('index.html', result=predicted_category, categories=categories, uploaded_text=content, filename=filename)
            except Exception as e:
                return render_template('index.html', result=f"Error: {str(e)}", categories=categories)
    return render_template('index.html', result=None, categories=categories)

@app.route('/feedback', methods=['POST'])
def feedback():
    try:
        feedback_text = request.form['feedback_text']
        is_correct = request.form['is_correct']
        feedback_label = request.form['feedback_label'] if is_correct == 'no' else None

        feedback_path = 'feedback.csv'
        feedback_data = {'text': [feedback_text], 'label': [feedback_label]}
        feedback_df = pd.DataFrame(feedback_data)
        if os.path.exists(feedback_path):
            feedback_df.to_csv(feedback_path, mode='a', header=False, index=False)
        else:
            feedback_df.to_csv(feedback_path, index=False)
        return "Feedback saved successfully!"
    except Exception as e:
        return f"Error saving feedback: {str(e)}", 500



if __name__ == '__main__':
    app.run(debug=True)
