import os
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import dill
import re
import nltk
from nltk.corpus import stopwords
from nltk.stem import PorterStemmer
from sklearn.feature_extraction.text import CountVectorizer
from sklearn.preprocessing import LabelEncoder
from sklearn.tree import DecisionTreeClassifier
from sklearn.model_selection import StratifiedKFold, cross_val_score
from sklearn.metrics import accuracy_score, classification_report, confusion_matrix
from scipy.sparse import vstack
from pptx import Presentation
import PyPDF2
from wordcloud import WordCloud
import csv

# Downloads stopwords if not done so
nltk.download('stopwords')

class DataLoader:
    """Load text data and labels from a directory and its subdirectories."""
    
    def __init__(self, directory):
        self.directory = directory

    def load_data(self):
        text_data, labels = [], []
        
        # Iterate through each subdirectory in the main directory
        for label in os.listdir(self.directory):
            label_dir = os.path.join(self.directory, label)
            if os.path.isdir(label_dir):  # Check if it's a directory
                for filename in os.listdir(label_dir):
                    file_path = os.path.join(label_dir, filename)
                    
                    if filename.endswith('.txt'):
                        with open(file_path, 'r', encoding='utf-8') as file:
                            text_data.append(file.read())
                        labels.append(label)

                    elif filename.endswith('.pdf'):
                        with open(file_path, 'rb') as file:
                            reader = PyPDF2.PdfReader(file)
                            content = ''.join(page.extract_text() for page in reader.pages)
                            text_data.append(content.strip())
                        labels.append(label)

                    elif filename.endswith('.pptx'):
                        presentation = Presentation(file_path)
                        content = ''.join(
                            shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")
                        )
                        text_data.append(content.strip())
                        labels.append(label)

        return np.array(text_data), np.array(labels)
    
class Preprocessor:
    """Preprocess text and encode labels."""

    def __init__(self, ngram_range=(1, 1)):
        self.vectorizer = CountVectorizer(ngram_range=ngram_range)
        self.label_encoder = LabelEncoder()
        self.stop_words = set(stopwords.words('english'))
        self.stemmer = PorterStemmer()

    def clean_text(self, text):
        return re.sub(r'[^a-zA-Z\s]', '', text.lower())

    def tokenize(self, text):
        return text.split()

    def remove_stopwords(self, tokens):
        return [word for word in tokens if word not in self.stop_words]

    def stem_tokens(self, tokens):
        return [self.stemmer.stem(word) for word in tokens]

    def preprocess_text(self, text_data):
        processed_data = []
        for text in text_data:
            tokens = self.stem_tokens(self.remove_stopwords(self.tokenize(self.clean_text(text))))
            processed_data.append(' '.join(tokens))
        return processed_data

    def count_vectorize(self, text_data):
        processed_data = self.preprocess_text(text_data)
        return self.vectorizer.fit_transform(processed_data)

    def encode_labels(self, labels):
        return self.label_encoder.fit_transform(labels)


class Model:

    def __init__(self, max_depth=None, min_samples_split=2):
        self.model = DecisionTreeClassifier(max_depth=max_depth, min_samples_split=min_samples_split)

    def train(self, X, y):
        # Dynamically adjust the number of splits to avoid errors with small classes but also ensuring at least 2 splits to avoid errors
        n_splits = max(2, min(5, np.min(np.bincount(y))))  # At least 2 splits, at most 5 or the smallest class size

        skf = StratifiedKFold(n_splits=n_splits, shuffle=True, random_state=42)

        scores = cross_val_score(self.model, X, y, cv=skf, scoring='accuracy')
        print(f'Cross-Validation Accuracy: {scores.mean():.2f} ± {scores.std():.2f}')
        # Fit the model on the entire dataset after cross-validation
        self.model.fit(X, y)

    def predict(self, X):
        return self.model.predict(X)

    def get_feature_importances(self):
        return self.model.feature_importances_


def main():
    directory_path = r'data'
    data_loader = DataLoader(directory_path)
    text_data, labels = data_loader.load_data()

    preprocessor = Preprocessor(ngram_range=(1, 2))
    X = preprocessor.count_vectorize(text_data)
    y = preprocessor.encode_labels(labels)

    model = Model(max_depth=10, min_samples_split=5)

    # Check for feedback data and update the model if available
    feedback_path = 'feedback.csv'
    if os.path.exists(feedback_path):
        feedback_data = pd.read_csv(feedback_path)
        feedback_text = feedback_data['text']
        feedback_labels = feedback_data['label']

        # Vectorize feedback
        
        feedback_X = preprocessor.vectorizer.transform(preprocessor.preprocess_text(feedback_text))
        feedback_y = preprocessor.label_encoder.transform(feedback_labels)


        if feedback_X.shape[0] == feedback_y.shape[0]:
            X = np.vstack([X, feedback_X])
            y = np.hstack([y, feedback_y])
            model.train(X, y)
        else:
            print("Feedback data dimensions do not match: Check feedback inputs.")

        # # Retrain model with feedback
        # X = np.vstack((X, feedback_X))
        # y = np.hstack((y, feedback_y))
        # model.train(X, y)

        # Save updated model
        with open('model.dill', 'wb') as f:
            dill.dump(model.model, f)

        
        # Clear feedback data
        open(feedback_path, 'w').close()  # This truncates the file to zero length
        print(f"Feedback data from {feedback_path} has been processed and cleared.")

        with open(feedback_path, mode='w', newline='\n', encoding='utf-8') as file:
            writer = csv.writer(file)
    
            writer.writerow(['text', 'label'])

            print(f"CSV file '{feedback_path}' has been initialised with headers 'text' and 'label'.")

main()